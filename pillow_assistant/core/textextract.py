"""Plain-text extraction for office docs (R2), shared by the preview panel and
the Agent's reference materialization. Optional deps degrade gracefully.
"""

from __future__ import annotations

from pathlib import Path

DOCX_MAX = 60_000


def read_docx_text(path: str | Path, max_chars: int = DOCX_MAX) -> str:
    """Extract paragraph + table text from a .docx. Requires python-docx."""
    import docx  # python-docx (optional)

    document = docx.Document(str(path))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    text = "\n".join(parts)
    if len(text) > max_chars:
        from pillow_assistant.core.i18n import t
        text = text[:max_chars] + t("extract.truncated")
    return text


def read_docx_html(path: str | Path) -> str:
    """Convert a .docx to formatted HTML for rich preview.

    Prefers mammoth (semantic HTML incl. headings/lists/tables/images); falls
    back to a basic HTML built from python-docx (headings, bold/italic, tables).
    """
    try:
        import mammoth  # optional, best fidelity

        with open(path, "rb") as fh:
            return mammoth.convert_to_html(fh).value
    except ImportError:
        return _docx_basic_html(path)


def _docx_basic_html(path: str | Path) -> str:
    import html as _html

    import docx  # python-docx

    document = docx.Document(str(path))
    parts = ["<html><body>"]
    for para in document.paragraphs:
        if not para.text.strip():
            parts.append("<br/>")
            continue
        style = (para.style.name or "").lower() if para.style else ""
        tag = "p"
        if style == "title" or "heading 1" in style:
            tag = "h1"
        elif "heading 2" in style:
            tag = "h2"
        elif "heading 3" in style:
            tag = "h3"
        inner = ""
        for run in para.runs:
            t = _html.escape(run.text)
            if run.bold:
                t = f"<b>{t}</b>"
            if run.italic:
                t = f"<i>{t}</i>"
            inner += t
        parts.append(f"<{tag}>{inner or _html.escape(para.text)}</{tag}>")
    for table in document.tables:
        parts.append("<table border='1' cellspacing='0' cellpadding='4'>")
        for row in table.rows:
            cells = "".join(f"<td>{_html.escape(c.text)}</td>" for c in row.cells)
            parts.append(f"<tr>{cells}</tr>")
        parts.append("</table>")
    parts.append("</body></html>")
    return "".join(parts)


def read_pptx_text(path: str | Path, max_chars: int = DOCX_MAX) -> str:
    """Extract per-slide text (shapes + tables) from a .pptx. Requires python-pptx."""
    from pptx import Presentation  # python-pptx (optional)

    prs = Presentation(str(path))
    parts: list[str] = []
    from pillow_assistant.core.i18n import t
    for i, slide in enumerate(prs.slides, 1):
        lines = [t("extract.slide", i=i)]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                lines.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append(" | ".join(c.text for c in row.cells))
        parts.append("\n".join(lines))
    text = "\n\n".join(parts)
    if len(text) > max_chars:
        text = text[:max_chars] + t("extract.truncated")
    return text
